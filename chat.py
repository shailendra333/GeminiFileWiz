import streamlit as st
from PyPDF2 import PdfReader
from docx import Document  # Word document reader
from pptx import Presentation  # PowerPoint reader
import pandas as pd  # Excel and CSV reader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_file_text(file):
    """Determines file type and extracts text accordingly."""
    if file.name.endswith('.pdf'):
        return extract_pdf_text(file)
    elif file.name.endswith('.docx'):
        return extract_word_text(file)
    elif file.name.endswith('.pptx'):
        return extract_ppt_text(file)
    elif file.name.endswith('.xlsx') or file.name.endswith('.xls'):
        return extract_excel_text(file)
    elif file.name.endswith('.csv'):
        return extract_csv_text(file)
    else:
        st.warning("Unsupported file format.")
        return ""

def extract_pdf_text(pdf_file):
    text = ""
    pdf_reader = PdfReader(pdf_file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_word_text(word_file):
    text = ""
    doc = Document(word_file)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_ppt_text(ppt_file):
    text = ""
    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_excel_text(excel_file):
    text = ""
    xls = pd.ExcelFile(excel_file)
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        text += df.to_string() + "\n"
    return text

def extract_csv_text(csv_file):
    df = pd.read_csv(csv_file)
    return df.to_string() + "\n"

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    # Updated prompt to align with StuffDocumentsChain expectations
    prompt_template = """
    Answer the question as detailed as possible from the provided context. 
    If the answer is not in the context, say, "Answer not available in the context."
    
    Context:
    {context}
    
    Question:
    {question}
    
    Answer:
    """
    
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    # Specify document_variable_name="context" explicitly
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt, document_variable_name="context")
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply:", response["output_text"])

def main():
    st.set_page_config(page_title="Chat with Documents")
    st.header("Engage with Your Documents via Gemini ✨📚📄")

    user_question = st.text_input("Ask a Question based on the Uploaded Documents")
    if user_question:
        user_input(user_question)

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your Documents (PDF, Word, PPT, Excel, CSV)", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                all_text = ""
                for file in uploaded_files:
                    all_text += get_file_text(file)
                if all_text:
                    text_chunks = get_text_chunks(all_text)
                    get_vector_store(text_chunks)
                    st.success("Document Processing Completed")

if __name__ == "__main__":
    main()
